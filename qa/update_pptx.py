"""
Actualiza el PowerPoint con:
  - Slide 15: Screenshot del dashboard real
  - Slide 17: Indicacion de videos demo disponibles
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import shutil

PPTX_SRC = r"C:\Users\cinth\Downloads\Monitoring IoT passif 247 Capteurs smartwatch + smartphone Score de risque 0-100 + alertes graduees (3 niveaux) Le medecin reste decideur final.pptx"
EVIDENCIAS = r"G:\Mi unidad\Cursos Master ADE\FIL ROUGE\Diagramas Drawio\mood-iot\qa\evidencias"
PPTX_OUT = r"C:\Users\cinth\Downloads\Mood-IoT_Presentation_Updated.pptx"

prs = Presentation(PPTX_SRC)
slide_w = prs.slide_width
slide_h = prs.slide_height

# ============================================================
# SLIDE 15 — Dashboard Medecin (agregar screenshot real)
# ============================================================
slide15 = prs.slides[14]

# Agregar screenshot del dashboard como imagen de fondo
dashboard_img = os.path.join(EVIDENCIAS, "screenshot_dashboard_hero.png")
if os.path.exists(dashboard_img):
    # Centrar la imagen en la parte inferior de la slide
    img_width = Inches(12.5)
    img_height = Inches(7.0)
    left = int((slide_w - img_width) / 2)
    top = Inches(2.2)
    pic = slide15.shapes.add_picture(dashboard_img, left, top, img_width, img_height)
    print(f"[OK] Slide 15: Dashboard screenshot added")

# ============================================================
# SLIDE 17 — Architecture - Demo (agregar info de videos)
# ============================================================
slide17 = prs.slides[16]

# Agregar screenshots del dashboard y mobile app
# Dashboard screenshot
if os.path.exists(dashboard_img):
    img_w = Inches(8)
    img_h = Inches(4.5)
    left = Inches(0.5)
    top = Inches(2.5)
    slide17.shapes.add_picture(dashboard_img, left, top, img_w, img_h)
    print(f"[OK] Slide 17: Dashboard screenshot added")

# Mobile screenshot
mobile_img = os.path.join(EVIDENCIAS, "sc_mob_05_sync_ok.png")
if os.path.exists(mobile_img):
    img_w = Inches(5.5)
    img_h = Inches(4.5)
    left = Inches(9)
    top = Inches(2.5)
    slide17.shapes.add_picture(mobile_img, left, top, img_w, img_h)
    print(f"[OK] Slide 17: Mobile screenshot added")

# Agregar texto indicando donde estan los videos
txBox = slide17.shapes.add_textbox(Inches(0.5), Inches(7.3), Inches(14), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Videos demo disponibles: demo_dashboard.webm + demo_mobile.webm"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
p.font.italic = True
p.alignment = PP_ALIGN.CENTER

# ============================================================
# SAVE
# ============================================================
prs.save(PPTX_OUT)
print(f"\n[OK] Presentacion actualizada guardada en:")
print(f"     {PPTX_OUT}")
print(f"     Videos demo en: {EVIDENCIAS}")
